"""
Build a PowerPoint summarising gravity-aligned Method 1 slope results for one run
(any --prefix: N75E_type3, sagamore, N75E, …).

Reads:   results/<prefix>_summary.json  and the per-flight PNGs that
         run_robinson.py wrote (<prefix>_<flight>_method1_profile.png / _segments.png).
Writes:  slides/<prefix>_report.pptx

    python build_type3_slides.py                 # default prefix N75E_type3
    python build_type3_slides.py sagamore        # any other run

Run it AFTER the analysis (ideally the gravity-aligned `--from-glb` pass) so the
numbers and plots it embeds are the corrected ones. Requires python-pptx
(`pip install python-pptx`).
"""
import os, sys, json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SLOPE_DIR = os.path.dirname(os.path.abspath(__file__))
RES = os.path.join(SLOPE_DIR, "results")
OUT_DIR = os.path.join(SLOPE_DIR, "slides")
os.makedirs(OUT_DIR, exist_ok=True)

PREFIX = sys.argv[1] if len(sys.argv) > 1 else "N75E_type3"

NAVY  = RGBColor(0x26, 0x46, 0x53)
TEAL  = RGBColor(0x2A, 0x9D, 0x8F)
RUST  = RGBColor(0xE7, 0x6F, 0x51)
GREY  = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK   = RGBColor(0x22, 0x22, 0x22)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ── helpers ───────────────────────────────────────────────────────────────────

def _txt(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    return tf

def band(slide, color=NAVY, h=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    s.shadow.inherit = False
    return s

def header(slide, title, sub=None, color=NAVY):
    band(slide, color, 1.15 if sub else 0.95)
    tf = _txt(slide, 0.5, 0.16, 12.3, 0.8)
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = WHITE
    if sub:
        tf2 = _txt(slide, 0.5, 0.98, 12.3, 0.45)
        p2 = tf2.paragraphs[0]; p2.text = sub
        p2.font.size = Pt(14); p2.font.color.rgb = RGBColor(0xDD, 0xE6, 0xE8)

def bullet(tf, text, size=16, bold=False, color=INK, first=False, level=0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = text; p.level = level
    p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color
    return p

def add_image_fit(slide, path, l, t, max_w, max_h):
    """Add image scaled to fit within (max_w, max_h) inches, keeping aspect."""
    if not (path and os.path.exists(path)):
        tf = _txt(slide, l, t + max_h / 2 - 0.3, max_w, 0.6)
        p = tf.paragraphs[0]; p.text = "[plot not found — run the analysis first]"
        p.font.size = Pt(12); p.font.italic = True; p.font.color.rgb = GREY
        p.alignment = PP_ALIGN.CENTER
        return
    from PIL import Image
    iw, ih = Image.open(path).size
    ar = iw / ih
    w, h = max_w, max_w / ar
    if h > max_h:
        h, w = max_h, max_h * ar
    l2 = l + (max_w - w) / 2
    t2 = t + (max_h - h) / 2
    slide.shapes.add_picture(path, Inches(l2), Inches(t2), Inches(w), Inches(h))

def fmt(v, suf="", nd=2):
    return f"{v:+.{nd}f}{suf}" if isinstance(v, (int, float)) else "—"


# ── load ──────────────────────────────────────────────────────────────────────

summary_path = os.path.join(RES, f"{PREFIX}_summary.json")
if not os.path.exists(summary_path):
    sys.exit(f"No summary at {summary_path} — run run_robinson.py --prefix {PREFIX} first.")
summary = json.load(open(summary_path))
if not summary:
    sys.exit(f"{summary_path} is empty.")

any_aligned = any(r.get("gravity_aligned") for r in summary)


# ── title slide ───────────────────────────────────────────────────────────────

s = prs.slides.add_slide(BLANK)
band(s, NAVY, SH)                       # full navy background
tf = _txt(s, 0.9, 2.4, 11.5, 1.6)
p = tf.paragraphs[0]; p.text = f"{PREFIX} — Slope Estimation"
p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = WHITE
p2 = tf.add_paragraph()
p2.text = ("VGGT reconstruction  →  Method 1 (ground-distance)  →  "
           "gravity-aligned to ground-truth camera poses")
p2.font.size = Pt(18); p2.font.color.rgb = TEAL
p3 = tf.add_paragraph()
p3.text = f"{len(summary)} flight(s):  " + ",  ".join(r["dataset"] for r in summary)
p3.font.size = Pt(14); p3.font.color.rgb = RGBColor(0xCC, 0xD5, 0xD7)


# ── method slide ──────────────────────────────────────────────────────────────

s = prs.slides.add_slide(BLANK)
header(s, "Method", "Why the estimate is trustworthy")
tf = _txt(s, 0.6, 1.35, 12.1, 5.7)
bullet(tf, "1.  VGGT reconstructs a 3-D point cloud + camera poses from the frames.",
       size=17, bold=True, color=NAVY, first=True)
bullet(tf, "2.  Method 1 reads the road surface beneath each camera and fits ground "
           "elevation vs. along-track distance — the line's tilt is the slope.", size=15)
bullet(tf, "3.  Gravity alignment (key fix):", size=17, bold=True, color=NAVY)
bullet(tf, "VGGT's reconstruction is only defined up to a similarity transform and is "
           "exported in camera-0's frame — its vertical axis is NOT gravity.", size=14, level=1)
bullet(tf, "So the raw slope MAGNITUDE is measured in a tilted frame and can be far off "
           "(a true 25° slope can read ~65°).", size=14, level=1, color=RUST)
bullet(tf, "Fix: align the VGGT camera centres to the ground-truth camera world positions "
           "(encoded in each frame's filename), recover the missing rotation, rotate the "
           "terrain into a true-vertical frame, then measure slope.", size=14, level=1, color=TEAL)
bullet(tf, "This corrects both the magnitude and the sign from one physical anchor. A small "
           "alignment residual (align_resid_m) confirms the fit is trustworthy.", size=14, level=1)


# ── per-flight slides ─────────────────────────────────────────────────────────

for r in summary:
    name = r["dataset"]
    final = r.get("final_signed_deg")
    direction = r.get("direction", "")
    s = prs.slides.add_slide(BLANK)
    hue = TEAL if (isinstance(final, (int, float)) and final >= 0) else RUST
    header(s, f"{name}", f"Final slope  {fmt(final,'°')}  {direction}", color=hue)

    # plots: profile (left) + segments (right)
    prof = os.path.join(RES, f"{PREFIX}_{name}_method1_profile.png")
    segs = os.path.join(RES, f"{PREFIX}_{name}_method1_segments.png")
    add_image_fit(s, prof, 0.4, 1.35, 6.2, 3.9)
    add_image_fit(s, segs, 6.8, 1.35, 6.2, 3.9)

    # stats strip
    tf = _txt(s, 0.5, 5.45, 12.3, 1.9)
    aligned = r.get("gravity_aligned")
    bullet(tf, f"Final (gravity-aligned): {fmt(final,'°')} {direction}"
               f"    |    Method 1 R² = {fmt(r.get('method1_r2'),'',3)}"
               f"    |    ground points = {r.get('n_ground_found','—')}",
           size=15, bold=True, color=NAVY, first=True)
    if aligned:
        bullet(tf, f"Gravity alignment: {r.get('align_source')} — residual "
                   f"{r.get('align_resid_m')} m over {r.get('align_frames')} frames  "
                   f"(raw VGGT-frame value was {fmt(r.get('method1_raw_vggt_deg'),'°')}, "
                   f"in a tilted frame).", size=13, color=GREY)
    else:
        bullet(tf, f"Gravity alignment NOT applied ({r.get('align_note')}). "
                   f"Falling back to: {r.get('estimate_basis')}. Magnitude in VGGT's "
                   f"tilted frame — treat with caution.", size=13, color=RUST)


# ── summary table slide ───────────────────────────────────────────────────────

s = prs.slides.add_slide(BLANK)
header(s, "Summary", "Gravity-aligned vs. raw VGGT-frame")
cols = ["Flight", "Final (aligned)", "Direction", "Raw VGGT", "R²", "Aligned?", "Resid (m)"]
rows = len(summary) + 1
tbl = s.shapes.add_table(rows, len(cols), Inches(0.6), Inches(1.5),
                         Inches(12.1), Inches(0.5 + 0.55 * len(summary))).table
for j, c in enumerate(cols):
    cell = tbl.cell(0, j); cell.text = c
    cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    para = cell.text_frame.paragraphs[0]
    para.font.size = Pt(13); para.font.bold = True; para.font.color.rgb = WHITE
for i, r in enumerate(summary, start=1):
    vals = [
        r["dataset"],
        fmt(r.get("final_signed_deg"), "°"),
        r.get("direction", "—"),
        fmt(r.get("method1_raw_vggt_deg"), "°"),
        fmt(r.get("method1_r2"), "", 3),
        "yes" if r.get("gravity_aligned") else "no",
        str(r.get("align_resid_m", "—")),
    ]
    for j, v in enumerate(vals):
        cell = tbl.cell(i, j); cell.text = str(v)
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(12); para.font.color.rgb = INK
        if j == 1:
            para.font.bold = True
            para.font.color.rgb = TEAL if str(v).startswith("+") else RUST

note = _txt(s, 0.6, 1.6 + 0.55 * len(summary) + 0.4, 12.1, 1.4)
bullet(note, "Final = gravity-aligned Method 1 (magnitude + sign from ground-truth poses). "
             "Raw VGGT = pre-alignment value in the tilted reconstruction frame, shown for "
             "transparency. A small residual means the pose alignment is reliable.",
       size=13, color=GREY, first=True)
if not any_aligned:
    bullet(note, "⚠  No flight in this run was gravity-aligned — numbers are raw VGGT-frame. "
                 "Re-run with the dataset path + flags so the poses can be read.",
           size=13, bold=True, color=RUST)


out = os.path.join(OUT_DIR, f"{PREFIX}_report.pptx")
prs.save(out)
print(f"Saved → {out}   ({len(prs.slides)} slides)")
