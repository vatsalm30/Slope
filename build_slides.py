"""
Build the slide deck (PowerPoint) summarizing Method 1 slope estimation:
method explanation, cross-application to the Method-2 dataset, per-segment
plots for every dataset, the flight-geometry finding, and the Robinson plan.

    python build_slides.py   ->   slides/method1_slope_report.pptx
"""
import os, json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SLOPE_DIR = os.path.dirname(os.path.abspath(__file__))
RES = os.path.join(SLOPE_DIR, "results")
OUT_DIR = os.path.join(SLOPE_DIR, "slides")
os.makedirs(OUT_DIR, exist_ok=True)

NAVY   = RGBColor(0x26, 0x46, 0x53)
TEAL   = RGBColor(0x2A, 0x9D, 0x8F)
RUST   = RGBColor(0xE7, 0x6F, 0x51)
GREY   = RGBColor(0x55, 0x55, 0x55)
LIGHT  = RGBColor(0xF4, 0xF1, 0xDE)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

summary = json.load(open(os.path.join(RES, "method1_summary.json")))
byname = {r["scenario"]: r for r in summary}


# ── helpers ───────────────────────────────────────────────────────────────────

def _txt(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    return tf

def band(slide, color=NAVY, h=1.0):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    s.shadow.inherit = False
    return s

def title_only(title, sub=None):
    s = prs.slides.add_slide(BLANK)
    band(s, NAVY, 1.15)
    tf = _txt(s, 0.5, 0.18, 12.3, 0.9)
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = WHITE
    if sub:
        tf2 = _txt(s, 0.5, 1.25, 12.3, 0.5)
        p2 = tf2.paragraphs[0]; p2.text = sub
        p2.font.size = Pt(15); p2.font.color.rgb = GREY
    return s

def bullets(tf, items, size=18, start=True):
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if (i == 0 and start) else tf.add_paragraph()
        lvl = 0; txt = it
        if isinstance(it, tuple):
            lvl, txt = it
        p.level = lvl
        p.text = ("• " if lvl == 0 else "– ") + txt
        p.font.size = Pt(size - 2 * lvl)
        p.font.color.rgb = NAVY if lvl == 0 else GREY
        p.space_after = Pt(6)

def pic_fit(slide, path, l, t, w, h):
    """Add picture scaled to fit box (l,t,w,h) inches, centered."""
    from PIL import Image
    iw, ih = Image.open(path).size
    box_w, box_h = w, h
    scale = min(box_w / iw, box_h / ih)
    pw, ph = iw * scale, ih * scale
    left = l + (box_w - pw) / 2
    top  = t + (box_h - ph) / 2
    slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(pw), Inches(ph))

def caption(slide, text, t, size=12):
    tf = _txt(slide, 0.5, t, 12.3, 0.4)
    p = tf.paragraphs[0]; p.text = text; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(size); p.font.italic = True; p.font.color.rgb = GREY


# ── 1. Title ────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
band(s, NAVY, 7.5)
tf = _txt(s, 0.8, 2.3, 11.7, 2.0)
p = tf.paragraphs[0]; p.text = "Method 1 (Ground-Distance) Slope Estimation"
p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = WHITE
p2 = tf.add_paragraph(); p2.text = "Cross-application to the Method-2 dataset · per-segment slope analysis · Robinson plan"
p2.font.size = Pt(18); p2.font.color.rgb = RGBColor(0xBF, 0xD7, 0xD2)
p3 = tf.add_paragraph(); p3.text = "Drone terrain slope from VGGT reconstructions"
p3.font.size = Pt(15); p3.font.color.rgb = RGBColor(0x9F, 0xBF, 0xBA)

# ── 2. Request + TL;DR ────────────────────────────────────────────────────────
s = title_only("What was asked, and the short answer")
tf = _txt(s, 0.6, 1.4, 6.1, 5.6)
bullets(tf, [
    "The request",
    (1, "Apply Method 1 to the dataset used by Method 2"),
    (1, "Method 1 is more general — it doesn't assume fixed height above ground"),
    (1, "Process the two Robinson datasets with Method 1"),
    (1, "Report slope with a sign (+ uphill / − downhill); use all images"),
    (1, "Per-segment slope plots; same plots for the previous dataset"),
    (1, "Collect everything into a slide file"),
], size=17)
tf2 = _txt(s, 6.95, 1.4, 5.8, 5.6)
bullets(tf2, [
    "Short answer",
    (1, "Method 1 transfers to the Method-2 dataset and works well"),
    (1, "Uphill: +3.7° (Method 2: +3.4°) — close agreement, R²=0.99"),
    (1, "Downhill: −7.9° (Method 2: −4.8°) — right direction, steeper"),
    (1, "Method 1 needs the reconstruction to capture vertical relief (terrain-following flights) — level flights rebuild flat"),
    (1, "Robinson uphill: +16.3° uphill (GPS-verified); downhill: −5.65° downhill"),
    (1, "Sign finding: Method 1's up/down sign rides on VGGT's vertical axis, which is NOT gravity-locked — we now anchor the sign to GPS"),
], size=17)

# ── 3. How Method 1 works ─────────────────────────────────────────────────────
s = title_only("How Method 1 works (ground-distance method)")
tf = _txt(s, 0.6, 1.35, 5.5, 5.6)
bullets(tf, [
    "Assume the GLB's up-axis (+Y) points along gravity",
    "For each camera, look straight down and take the road surface beneath it (point cloud)",
    "Collect those ground points along the whole flight",
    "Fit a line: ground elevation vs. distance along the path",
    "The tilt of that line is the slope",
    "It reads the terrain directly — so it does NOT need the drone to keep a fixed height (that is Method 2's assumption)",
], size=18)
pic_fit(s, os.path.join(RES, "fixed_distance2ground_downhill_profile.png"), 6.2, 1.5, 6.7, 4.8)
caption(s, "The actual fit: ground elevation vs. along-track distance — its tilt is the slope.", 6.45)

# ── 4. Method 1 vs Method 2 ───────────────────────────────────────────────────
s = title_only("Method 1 vs. Method 2 — what each assumes")
rows = [
    ["", "Method 1 — ground-distance", "Method 2 — camera-altitude"],
    ["What it measures", "Terrain points beneath the camera", "The camera's own up/down motion"],
    ["Key assumption", "Up-axis = gravity", "Drone keeps a fixed height above ground"],
    ["Generality", "Works for any flight (if relief is reconstructed)", "Only valid for terrain-following flights"],
    ["Fails when", "Reconstruction comes back flat (no vertical parallax)", "Drone altitude ≠ terrain (e.g. level flight)"],
]
tbl = s.shapes.add_table(len(rows), 3, Inches(0.6), Inches(1.5),
                         Inches(12.1), Inches(3.8)).table
tbl.columns[0].width = Inches(2.7)
tbl.columns[1].width = Inches(4.7)
tbl.columns[2].width = Inches(4.7)
for r, row in enumerate(rows):
    for c, val in enumerate(row):
        cell = tbl.cell(r, c); cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14 if r else 15)
            p.font.bold = (r == 0 or c == 0)
            p.font.color.rgb = WHITE if r == 0 else NAVY
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY if r == 0 else (RGBColor(0xED,0xF3,0xF2) if r % 2 else WHITE)
tf = _txt(s, 0.6, 5.6, 12.1, 1.2)
bullets(tf, [
    "Shengqing's point is right: Method 1 is the more general of the two — it never relies on the fixed-height assumption.",
], size=16)

# ── 5. Sign convention ────────────────────────────────────────────────────────
s = title_only("Sign convention (so + / − is unambiguous)")
tf = _txt(s, 0.7, 1.6, 11.9, 4.0)
bullets(tf, [
    "All images are used (1, 2, 3, …) — no frames dropped.",
    "The sign is anchored to image order — the direction the drone actually flew:",
    (1, "+  =  terrain goes UPHILL as the flight progresses"),
    (1, "−  =  terrain goes DOWNHILL as the flight progresses"),
    "Each dataset reports one overall slope (a single line fit over all images) plus a per-segment breakdown between consecutive images.",
    "R² is shown alongside: high R² = the terrain profile is cleanly reconstructed; near-zero R² = no usable slope signal.",
], size=18)

# ── 6. Results table ──────────────────────────────────────────────────────────
s = title_only("Results — all datasets (Method 1, all images)")
order = ["happy_hollow_11views", "fixed_distance2ground_uphill", "fixed_distance2ground_downhill",
         "fixed_altitude_uphill", "fixed_altitude_downhill"]
labels = {
    "happy_hollow_11views": "Happy Hollow (original)",
    "fixed_distance2ground_uphill": "fixed_dist2ground — uphill",
    "fixed_distance2ground_downhill": "fixed_dist2ground — downhill",
    "fixed_altitude_uphill": "fixed_altitude — uphill",
    "fixed_altitude_downhill": "fixed_altitude — downhill",
}
hdr = ["Dataset", "Method 1", "R²", "Method 2", "Verdict"]
rows = [hdr]
for k in order:
    if k not in byname: continue
    r = byname[k]
    m2 = f"{r['method2_signed']:+.2f}°" if r["method2_signed"] is not None else "—"
    verdict = "clean signal" if r["method1_r2"] > 0.5 else "no signal (flat rebuild)"
    rows.append([labels[k], f"{r['method1_signed']:+.2f}°", f"{r['method1_r2']:.3f}", m2, verdict])
tbl = s.shapes.add_table(len(rows), 5, Inches(0.6), Inches(1.6),
                         Inches(12.1), Inches(3.4)).table
widths = [4.1, 1.9, 1.4, 1.9, 2.8]
for i, w in enumerate(widths): tbl.columns[i].width = Inches(w)
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci); cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14 if ri else 15); p.font.bold = (ri == 0)
            p.font.color.rgb = WHITE if ri == 0 else NAVY
            if ci == 4 and ri > 0:
                p.font.color.rgb = TEAL if "clean" in val else RUST
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY if ri == 0 else (RGBColor(0xED,0xF3,0xF2) if ri % 2 else WHITE)
caption(s, "Method 2 numbers shown only for the fixed_dist2ground datasets (where Method 2 is valid).", 5.2)

# ── 7. Cross-application: Method 1 on Method-2 dataset ────────────────────────
s = title_only("Method 1 on the Method-2 dataset (the main request)",
               "fixed_distance2ground — terrain-following flight")
pic_fit(s, os.path.join(RES, "fixed_distance2ground_uphill_profile.png"), 0.4, 1.55, 6.3, 3.4)
pic_fit(s, os.path.join(RES, "fixed_distance2ground_downhill_profile.png"), 6.7, 1.55, 6.3, 3.4)
tf = _txt(s, 0.6, 5.15, 12.1, 2.1)
bullets(tf, [
    "Uphill +3.69° (R²=0.99) vs Method 2 +3.44°  →  excellent agreement.",
    "Downhill −7.90° (R²=0.76) vs Method 2 −4.80°  →  correct direction, magnitude steeper.",
    "Takeaway: Method 1 transfers cleanly here because this flight follows the terrain, giving VGGT vertical parallax to reconstruct the slope.",
], size=16)

# ── 8. Per-segment plots for the Method-2 dataset ─────────────────────────────
s = title_only("Per-segment slopes — Method-2 dataset (fixed_distance2ground)")
pic_fit(s, os.path.join(RES, "fixed_distance2ground_uphill_method1_segments.png"), 0.3, 1.5, 6.4, 4.6)
pic_fit(s, os.path.join(RES, "fixed_distance2ground_downhill_method1_segments.png"), 6.6, 1.5, 6.4, 4.6)
caption(s, "x = image-index segment, y = signed slope. Single near-vertical bars (closely-spaced cameras) are capped & flagged.", 6.3)

# ── 9. Previous dataset (fixed_altitude) ──────────────────────────────────────
s = title_only("Previous dataset — Method 1 on fixed_altitude (level flight)")
pic_fit(s, os.path.join(RES, "fixed_altitude_uphill_method1_segments.png"), 0.3, 1.5, 6.4, 4.3)
pic_fit(s, os.path.join(RES, "fixed_altitude_downhill_method1_segments.png"), 6.6, 1.5, 6.4, 4.3)
tf = _txt(s, 0.6, 5.95, 12.1, 1.4)
bullets(tf, [
    "Both come out near 0° with R²≈0.00 — the per-segment slopes are just noise (alternating up/down).",
    "This is NOT terrain being flat: it is the SAME hill as the fixed_distance2ground flights (matching GPS track).",
], size=15)

# ── 10. The finding: flight geometry ──────────────────────────────────────────
s = title_only("Why: the flight geometry decides whether Method 1 has a signal",
               "Same terrain, same GPS path, two flight styles")
pic_fit(s, os.path.join(RES, "fixed_distance2ground_downhill_profile.png"), 0.4, 1.55, 6.3, 3.7)
pic_fit(s, os.path.join(RES, "fixed_altitude_downhill_profile.png"), 6.7, 1.55, 6.3, 3.7)
tf = _txt(s, 0.6, 5.4, 12.1, 1.9)
bullets(tf, [
    "Left (terrain-following): camera moves up/down with the hill → VGGT reconstructs the relief → clean −7.9° fit.",
    "Right (level / fixed-altitude): camera moves horizontally only → little vertical parallax → ground rebuilds flat → no slope to fit.",
    "So 'up-axis = gravity' is necessary but not sufficient — the reconstruction must also preserve vertical relief.",
], size=15)

# ── 11. Happy Hollow ──────────────────────────────────────────────────────────
s = title_only("Reference: original Happy Hollow reconstruction (Method 1)")
pic_fit(s, os.path.join(RES, "happy_hollow_profile.png"), 0.4, 1.55, 6.3, 4.6)
pic_fit(s, os.path.join(RES, "happy_hollow_method1_segments.png"), 6.7, 1.6, 6.3, 4.4)
caption(s, "Original 11-view reconstruction: −3.92° downhill, R²=0.99 — Method 1's clean reference case.", 6.35)

# ── 12. The sign problem + GPS fix ────────────────────────────────────────────
s = title_only("Fixing the up/down sign — anchor to GPS, not VGGT",
               "VGGT has no gravity sensor, so its vertical axis (and the slope sign) is not reliable")
tf = _txt(s, 0.6, 1.35, 6.0, 5.7)
bullets(tf, [
    "Method 1 reads slope in VGGT's coordinate frame, assuming +Y = gravity.",
    "VGGT never measures gravity — it picks an arbitrary frame — so the up/down SIGN can come out backwards.",
    "Proof on Robinson uphill: the photos carry GPS, and the drone climbed 168.7 → 172.3 m (clearly uphill).",
    (1, "GPS:      +16.3° uphill  (R²=0.93)"),
    (1, "Method 1: −2.6°  downhill (R²=0.29)  ✗ wrong sign & noisy"),
    "Fix: when photos carry GPS, take the sign (and a metric slope) from the GPS track. Otherwise fall back to the known flight label.",
], size=16)
prof_gps = os.path.join(RES, "robinson_uphill_gps_track.png")
if os.path.exists(prof_gps):
    pic_fit(s, prof_gps, 6.7, 1.5, 6.3, 4.9)
    caption(s, "Robinson uphill: GPS altitude vs. along-track distance — the gravity-true profile.", 6.5)

# ── 13. Robinson — results ────────────────────────────────────────────────────
robinson_json = os.path.join(RES, "robinson_summary.json")
robinson = json.load(open(robinson_json)) if os.path.exists(robinson_json) else []

if robinson:
    s = title_only("Robinson datasets — results", "all images, signed (+ uphill / − downhill)")
    hdr = ["Dataset", "Imgs", "Method 1 (VGGT)", "GPS track", "FINAL", "Sign from"]
    rows = [hdr]
    for r in robinson:
        m1c = f"{r['method1_signed_deg']:+.2f}° (R²={r['method1_r2']:.2f})"
        gpsc = (f"{r['gps_signed_deg']:+.2f}° (R²={r['gps_r2']:.2f})"
                if r.get("gps_available") else "—")
        finalc = f"{r['final_signed_deg']:+.2f}° {r['direction']}"
        rows.append([r["dataset"], str(r["n_images"]), m1c, gpsc, finalc, r["sign_source"]])
    tbl = s.shapes.add_table(len(rows), 6, Inches(0.4), Inches(1.7),
                             Inches(12.5), Inches(2.4)).table
    for i, w in enumerate([1.9, 0.9, 3.1, 2.7, 2.3, 1.6]): tbl.columns[i].width = Inches(w)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci); cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12 if ri else 12); p.font.bold = (ri == 0 or ci == 4)
                p.font.color.rgb = WHITE if ri == 0 else NAVY
                if ci == 4 and ri > 0:
                    p.font.color.rgb = TEAL if val.strip().startswith("+") else RUST
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if ri == 0 else (RGBColor(0xED,0xF3,0xF2) if ri % 2 else WHITE)
    tf = _txt(s, 0.5, 4.5, 12.4, 2.6)
    bullets(tf, [
        "uphill — iPhone 14 with GPS: drone climbed; FINAL +16.3° uphill from the GPS track (Method 1's ground-fit was too noisy here, R²=0.29).",
        "downhill — DJI, GPS stripped from the files: Method 1's ground-fit is clean (R²≈1.00), so we keep its 5.65° magnitude and take the sign from the folder label → −5.65° downhill.",
        "To verify downhill independently, we'd need its flight log / original GPS — flagged for follow-up.",
    ], size=14)
    caption(s, "FINAL sign source: GPS where available, otherwise the known flight label. Magnitude: GPS when Method 1's R² is low, else Method 1.", 7.0, size=11)

    # One detail slide per Robinson dataset
    for r in robinson:
        nm = r["dataset"]
        gps_png = os.path.join(RES, f"robinson_{nm}_gps_track.png")
        prof = os.path.join(RES, f"robinson_{nm}_method1_profile.png")
        seg  = os.path.join(RES, f"robinson_{nm}_method1_segments.png")
        s = title_only(f"Robinson — {nm}",
                       f"FINAL {r['final_signed_deg']:+.2f}° {r['direction']}   "
                       f"[{r['estimate_basis']}]   ·   {r['n_images']} images")
        # left: the trusted basis (GPS track if present, else Method 1 profile)
        left = gps_png if (r.get("gps_available") and os.path.exists(gps_png)) else prof
        if left and os.path.exists(left): pic_fit(s, left, 0.3, 1.55, 6.4, 4.6)
        if os.path.exists(seg):           pic_fit(s, seg, 6.6, 1.55, 6.4, 4.6)
        lcap = ("Left: GPS track (gravity-true)." if r.get("gps_available")
                else "Left: Method 1 ground-elevation fit.")
        caption(s, f"{lcap}  Right: Method 1 per-segment slope by image index.", 6.35)
else:
    s = title_only("Robinson datasets — run on the A100 box")
    tf = _txt(s, 0.7, 1.5, 11.9, 5.2)
    bullets(tf, [
        "python run_robinson.py \"robinson copy\" --multi      (GPU: reconstruct + analyze)",
        "python run_robinson.py --from-glb                      (re-analyze existing GLBs, no GPU)",
        "Each dataset reports Method 1 (VGGT), the GPS-track slope, and a GPS/label-anchored FINAL sign.",
    ], size=17)

# ── 13. Observations ──────────────────────────────────────────────────────────
s = title_only("Observations & recommendations")
tf = _txt(s, 0.7, 1.5, 11.9, 5.2)
bullets(tf, [
    "Method 1 estimates slope MAGNITUDE well when R² is high — but its up/down SIGN is unreliable, because VGGT's vertical axis is not gravity-locked.",
    "Always anchor the sign to a gravity-true reference: GPS altitude is ideal (we proved it flips the wrong Robinson-uphill sign to the correct +16.3°).",
    "Capture GPS/flight-log with every dataset — without it, the sign can only be taken from the known flight label, which we can't independently verify (e.g. Robinson downhill).",
    "Reliability tracks R²: trust the magnitude when R² is high; treat near-zero R² as 'no signal', not 'flat ground'.",
    "Per-segment slopes are noisy and can spike — read them as a profile, not point values; use all images for the overall fit.",
    "Prefer flights with vertical parallax (terrain-following / oblique) over level passes.",
], size=16)

# ── 14. Reproduce ─────────────────────────────────────────────────────────────
s = title_only("Files & how to reproduce")
tf = _txt(s, 0.7, 1.5, 11.9, 5.2)
bullets(tf, [
    "method1_slope.py — Method 1 (all-frames, signed, image-order, adaptive ground-search radius) + plot helpers",
    "gps_anchor.py — reads EXIF GPS, computes the gravity-true track slope; the sign anchor",
    "run_robinson.py — A100: reconstruct with VGGT (or --from-glb to skip GPU), then Method 1 + GPS anchor",
    "run_method1_existing.py — Method 1 on the earlier GLBs → results/ plots + summary.json",
    "build_slides.py — regenerates this deck from results/",
    "Re-run order:  (A100) run_robinson.py [--from-glb]  →  build_slides.py",
], size=17)

out = os.path.join(OUT_DIR, "method1_slope_report.pptx")
prs.save(out)
print(f"Saved → {out}  ({len(prs.slides._sldIdLst)} slides)")
